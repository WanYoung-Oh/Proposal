from pathlib import Path

from pptx import Presentation

from .base import ParsedDocument, ParsedSlide


def _extract_text(shape) -> str:
    if not shape.has_text_frame:
        return ""
    lines = []
    for para in shape.text_frame.paragraphs:
        line = "".join(run.text for run in para.runs).strip()
        if line:
            lines.append(line)
    return "\n".join(lines)


def _is_title_placeholder(shape) -> bool:
    try:
        ph = shape.placeholder_format
        return ph is not None and ph.idx == 0
    except Exception:
        return False


def parse_pptx(path: Path, doc_id: str) -> ParsedDocument:
    prs = Presentation(str(path))
    slides = []

    for i, slide in enumerate(prs.slides, 1):
        title_text = ""
        body_parts = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if _is_title_placeholder(shape):
                title_text = shape.text_frame.text.strip()
            else:
                text = _extract_text(shape)
                if text:
                    body_parts.append(text)

        notes_text = ""
        if slide.has_notes_slide:
            tf = slide.notes_slide.notes_text_frame
            notes_text = tf.text.strip() if tf else ""

        slides.append(ParsedSlide(
            slide_no=i,
            title=title_text,
            body="\n".join(body_parts),
            notes=notes_text,
        ))

    return ParsedDocument(
        doc_id=doc_id,
        source_path=str(path),
        file_type="pptx",
        slides=slides,
    )

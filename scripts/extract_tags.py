#!/usr/bin/env python3
"""
제안서/RFP에서 평가 항목, 전략 키워드, 핵심 차별화 포인트 추출
→ 각 프로젝트 폴더에 tags.json 저장
"""

import json
import os
import sys
import time
from pathlib import Path

import anthropic
import pdfplumber
from pptx import Presentation

PROJECTS_DIR = Path(__file__).parent.parent / "data" / "projects"

MAX_CHARS = 40000  # Claude에 전달할 최대 텍스트 길이


# ─────────────────────────────────────────────────────────────────
# 텍스트 추출
# ─────────────────────────────────────────────────────────────────

def extract_pdf(path: Path) -> str:
    texts = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                texts.append(t)
    return "\n".join(texts)


def extract_pptx(path: Path) -> str:
    import zipfile
    try:
        prs = Presentation(path)
    except zipfile.BadZipFile:
        print(f"  ⚠️  손상된 PPTX 파일, 건너뜀: {path.name}")
        return ""
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(r.text for r in para.runs if r.text.strip())
                    if line.strip():
                        texts.append(line.strip())
    return "\n".join(texts)


def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_pdf(path)
    elif suffix in (".pptx", ".ppt"):
        return extract_pptx(path)
    return ""


# ─────────────────────────────────────────────────────────────────
# Claude API 호출
# ─────────────────────────────────────────────────────────────────

RFP_PROMPT = """다음은 공공정보화 사업의 제안요청서(RFP) 텍스트입니다.

아래 항목을 JSON 형식으로 추출해 주세요:
- evaluation_criteria: 평가 항목명과 배점 목록 (없으면 빈 배열)

반드시 아래 JSON 형식만 출력하세요 (설명 없이):
{{
  "evaluation_criteria": [
    {{"item": "항목명", "score": 숫자}},
    ...
  ]
}}

텍스트:
{text}
"""

PROPOSAL_PROMPT = """다음은 공공정보화 사업의 제안서 텍스트입니다.

아래 항목을 JSON 형식으로 추출해 주세요:
- strategy_keywords: 제안사가 강조하는 핵심 전략 키워드 목록 (최대 15개, 짧은 명사구)
- differentiators: 경쟁사 대비 핵심 차별화 포인트 목록 (최대 10개, 한 문장)
- strategy_summary: 제안사의 핵심 제안전략 방향 3~5개. 각 항목은 "무엇을 어떻게 하여 어떤 효과를 달성한다" 구조의 완성된 문장으로 작성. RAG 컨텍스트로 활용되므로 구체적이고 독립적으로 이해 가능하게 작성.

반드시 아래 JSON 형식만 출력하세요 (설명 없이):
{{
  "strategy_keywords": ["키워드1", "키워드2", ...],
  "differentiators": ["차별화 포인트1", ...],
  "strategy_summary": ["전략 방향 문장1", "전략 방향 문장2", ...]
}}

텍스트:
{text}
"""

STRATEGY_SUMMARY_PROMPT = """다음은 공공정보화 사업의 제안서 텍스트입니다.

제안사의 핵심 제안전략 방향을 3~5개의 완성된 문장으로 요약해 주세요.
각 문장은 "무엇을 어떻게 하여 어떤 효과를 달성한다" 구조로 구체적으로 작성합니다.
RAG 컨텍스트로 활용되므로 맥락 없이도 독립적으로 이해 가능해야 합니다.

반드시 아래 JSON 형식만 출력하세요 (설명 없이):
{{
  "strategy_summary": ["전략 방향 문장1", "전략 방향 문장2", ...]
}}

텍스트:
{text}
"""


def call_claude(client: anthropic.Anthropic, prompt: str, max_tokens: int = 2048) -> dict:
    response = client.messages.create(
        model="claude-haiku-4-5-20251001",
        max_tokens=max_tokens,
        messages=[{"role": "user", "content": prompt}],
    )
    raw = response.content[0].text.strip()
    # JSON 블록 추출
    if "```" in raw:
        raw = raw.split("```")[1]
        if raw.startswith("json"):
            raw = raw[4:]
    return json.loads(raw)


# ─────────────────────────────────────────────────────────────────
# 메인
# ─────────────────────────────────────────────────────────────────

def _load_proposal_text(proj_dir: Path) -> str:
    for ext in (".pdf", ".pptx", ".ppt"):
        path = proj_dir / f"proposal{ext}"
        if path.exists():
            return extract_text(path)[:MAX_CHARS]
    return ""


def extract_strategy_summary(client: anthropic.Anthropic, proj_dir: Path) -> list[str]:
    """strategy_summary만 추출 (기존 tags.json 부분 업데이트용)"""
    text = _load_proposal_text(proj_dir)
    if not text.strip():
        return []
    try:
        result = call_claude(client, STRATEGY_SUMMARY_PROMPT.format(text=text))
        return result.get("strategy_summary", [])
    except Exception as e:
        print(f"  ⚠️  전략 요약 추출 실패: {e}")
        return []


def process_project(client: anthropic.Anthropic, proj_dir: Path) -> dict:
    tags = {
        "evaluation_criteria": [],
        "strategy_keywords": [],
        "differentiators": [],
        "strategy_summary": [],
    }

    # RFP → 평가 항목
    rfp_path = proj_dir / "rfp.pdf"
    if rfp_path.exists():
        text = extract_text(rfp_path)[:MAX_CHARS]
        if text.strip():
            try:
                result = call_claude(client, RFP_PROMPT.format(text=text))
                tags["evaluation_criteria"] = result.get("evaluation_criteria", [])
            except Exception as e:
                print(f"  ⚠️  RFP 추출 실패: {e}")

    # Proposal → 전략 키워드 + 차별화 포인트 + 전략 요약 (단일 호출)
    text = _load_proposal_text(proj_dir)
    if text.strip():
        try:
            result = call_claude(client, PROPOSAL_PROMPT.format(text=text))
            tags["strategy_keywords"] = result.get("strategy_keywords", [])
            tags["differentiators"]   = result.get("differentiators", [])
            tags["strategy_summary"]  = result.get("strategy_summary", [])
        except Exception as e:
            print(f"  ⚠️  Proposal 추출 실패: {e}")

    return tags


def main():
    api_key = os.environ.get("ANTHROPIC_API_KEY")
    if not api_key:
        print("❌ ANTHROPIC_API_KEY 환경변수가 설정되지 않았습니다.")
        sys.exit(1)

    client = anthropic.Anthropic(api_key=api_key)
    proj_dirs = sorted(p for p in PROJECTS_DIR.iterdir() if p.is_dir())

    # 특정 프로젝트만 실행할 경우 인자로 전달
    if len(sys.argv) > 1:
        proj_dirs = [p for p in proj_dirs if p.name in sys.argv[1:]]

    print(f"총 {len(proj_dirs)}개 프로젝트 처리\n")

    for i, proj_dir in enumerate(proj_dirs, 1):
        pid = proj_dir.name
        tags_path = proj_dir / "tags.json"

        if tags_path.exists():
            existing = json.loads(tags_path.read_text(encoding="utf-8"))

            if "strategy_summary" in existing:
                print(f"[{i:02d}/{len(proj_dirs)}] 스킵 (완료): {pid}")
                continue

            # strategy_summary만 추가 — evaluation_criteria 등 기존 필드 보존
            print(f"[{i:02d}/{len(proj_dirs)}] strategy_summary 추가 중: {pid}")
            summary = extract_strategy_summary(client, proj_dir)
            existing["strategy_summary"] = summary
            tags_path.write_text(json.dumps(existing, ensure_ascii=False, indent=2), encoding="utf-8")
            print(f"  전략 요약: {len(summary)}개  (기존 필드 보존)")

        else:
            print(f"[{i:02d}/{len(proj_dirs)}] 신규 처리 중: {pid}")
            tags = process_project(client, proj_dir)
            tags_path.write_text(json.dumps(tags, ensure_ascii=False, indent=2), encoding="utf-8")
            print(f"  평가항목: {len(tags['evaluation_criteria'])}개 "
                  f"| 전략키워드: {len(tags['strategy_keywords'])}개 "
                  f"| 차별화: {len(tags['differentiators'])}개 "
                  f"| 전략요약: {len(tags['strategy_summary'])}개")

        time.sleep(0.3)  # API rate limit 방지

    print("\n✅ 완료")


if __name__ == "__main__":
    main()
